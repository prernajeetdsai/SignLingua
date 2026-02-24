"""
main.py — SignLingua Core Backend  (v3 — Enhanced)
───────────────────────────────────────────────────
STT  : OpenAI Whisper  — local, no API key, supports 99 languages
       Input language hint passed to Whisper for accuracy boost
OCR  : Tesseract  — multi-PSM, image preprocessing pipeline
Trans: Google Translate via deep-translator
TTS  : gTTS  — 40+ languages
AI   : Gemini 2.5 Flash (google-genai SDK)  +  FAISS RAG embeddings
"""

import os, io, tempfile, textwrap
from typing import Optional


# ════════════════════════════════════════════════════════════════
# SUPPORTED LANGUAGES
# ════════════════════════════════════════════════════════════════

SUPPORTED_LANGUAGES: dict[str, str] = {
    "Auto-Detect":            "auto",
    "English":                "en",
    "Hindi":                  "hi",
    "French":                 "fr",
    "Spanish":                "es",
    "German":                 "de",
    "Chinese (Simplified)":   "zh-cn",
    "Chinese (Traditional)":  "zh-tw",
    "Japanese":               "ja",
    "Korean":                 "ko",
    "Arabic":                 "ar",
    "Portuguese":             "pt",
    "Russian":                "ru",
    "Italian":                "it",
    "Dutch":                  "nl",
    "Turkish":                "tr",
    "Bengali":                "bn",
    "Tamil":                  "ta",
    "Telugu":                 "te",
    "Marathi":                "mr",
    "Gujarati":               "gu",
    "Kannada":                "kn",
    "Malayalam":              "ml",
    "Punjabi":                "pa",
    "Urdu":                   "ur",
    "Thai":                   "th",
    "Vietnamese":             "vi",
    "Indonesian":             "id",
    "Malay":                  "ms",
    "Greek":                  "el",
    "Polish":                 "pl",
    "Swedish":                "sv",
    "Norwegian":              "no",
    "Danish":                 "da",
    "Finnish":                "fi",
    "Czech":                  "cs",
    "Romanian":               "ro",
    "Hungarian":              "hu",
    "Ukrainian":              "uk",
    "Hebrew":                 "iw",
    "Swahili":                "sw",
    "Afrikaans":              "af",
}

# For target language selector (exclude Auto-Detect)
TARGET_LANGUAGES: dict[str, str] = {
    k: v for k, v in SUPPORTED_LANGUAGES.items() if v != "auto"
}

LANG_CODE_TO_NAME: dict[str, str] = {v: k for k, v in SUPPORTED_LANGUAGES.items()}
RTL_LANGS = {"ar", "he", "ur", "iw", "fa"}

# Whisper ISO-639-1 code mapping (Whisper uses full names for some)
WHISPER_LANG_MAP: dict[str, str] = {
    "auto":    None,       # None = Whisper auto-detects
    "zh-cn":   "zh",
    "zh-tw":   "zh",
    "iw":      "he",       # Whisper uses 'he' for Hebrew
}

# Whisper model — lazy loaded once
_whisper_model = None


# ════════════════════════════════════════════════════════════════
# OCR
# ════════════════════════════════════════════════════════════════

def extract_text_from_image(image_bytes: bytes, hint_lang: str = "auto") -> str:
    """
    Tesseract OCR with multi-PSM strategy and image preprocessing.
    hint_lang: ISO code to help Tesseract with language-specific characters.
    """
    try:
        import pytesseract
        from PIL import Image, ImageEnhance, ImageFilter

        img = Image.open(io.BytesIO(image_bytes))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")

        # Preprocessing pipeline for better accuracy
        img = ImageEnhance.Contrast(img).enhance(1.8)
        img = ImageEnhance.Sharpness(img).enhance(2.2)
        img = img.filter(ImageFilter.MedianFilter(size=3))

        # Map our lang code to Tesseract language pack
        tess_lang = _to_tesseract_lang(hint_lang)
        lang_arg = f"-l {tess_lang}" if tess_lang else ""

        best = ""
        for psm in ("--psm 6", "--psm 11", "--psm 3", "--psm 4"):
            cfg = f"{lang_arg} {psm}".strip()
            try:
                t = pytesseract.image_to_string(img, config=cfg).strip()
                if len(t) > len(best):
                    best = t
            except Exception:
                continue

        if not best:
            raise ValueError("No text could be extracted from the image.")
        return best

    except ValueError:
        raise
    except Exception as exc:
        raise RuntimeError(f"OCR failed: {exc}") from exc


def _to_tesseract_lang(lang_code: str) -> str:
    """Map ISO 639-1 code → Tesseract language pack name."""
    mapping = {
        "en": "eng", "hi": "hin", "fr": "fra", "es": "spa", "de": "deu",
        "zh-cn": "chi_sim", "zh-tw": "chi_tra", "ja": "jpn", "ko": "kor",
        "ar": "ara", "pt": "por", "ru": "rus", "it": "ita", "nl": "nld",
        "tr": "tur", "bn": "ben", "ta": "tam", "te": "tel", "mr": "mar",
        "gu": "guj", "kn": "kan", "ml": "mal", "pa": "pan", "ur": "urd",
        "th": "tha", "vi": "vie", "el": "ell", "pl": "pol", "uk": "ukr",
        "iw": "heb", "sw": "swa", "af": "afr",
    }
    return mapping.get(lang_code, "eng")


# ════════════════════════════════════════════════════════════════
# LANGUAGE DETECTION
# ════════════════════════════════════════════════════════════════

def detect_language(text: str) -> str:
    """Return ISO-639-1 code or 'unknown'."""
    try:
        from langdetect import detect, DetectorFactory
        DetectorFactory.seed = 42
        return detect(text)
    except Exception:
        return "unknown"


# ════════════════════════════════════════════════════════════════
# TRANSLATION
# ════════════════════════════════════════════════════════════════

def translate_text(text: str, target: str, source: str = "auto") -> dict:
    """Google Translate via deep-translator."""
    try:
        from deep_translator import GoogleTranslator

        # deep-translator accepts "auto" for source
        src = source if source != "auto" else "auto"
        translated = GoogleTranslator(source=src, target=target).translate(text)
        detected_src = detect_language(text) if source == "auto" else source

        return {
            "translated_text": translated,
            "source_language": detected_src,
            "target_language": target,
        }
    except Exception as exc:
        raise RuntimeError(f"Translation failed: {exc}") from exc


# ════════════════════════════════════════════════════════════════
# TEXT-TO-SPEECH
# ════════════════════════════════════════════════════════════════

def text_to_speech(text: str, lang_code: str) -> bytes:
    """gTTS → MP3 bytes. Returns None if lang unsupported."""
    try:
        from gtts import gTTS

        # gTTS quirks
        code_map = {"zh-cn": "zh-CN", "zh-tw": "zh-TW", "iw": "iw"}
        code = code_map.get(lang_code, lang_code)

        buf = io.BytesIO()
        gTTS(text=text, lang=code, slow=False).write_to_fp(buf)
        buf.seek(0)
        return buf.read()
    except Exception as exc:
        raise RuntimeError(f"TTS failed: {exc}") from exc


# ════════════════════════════════════════════════════════════════
# FILE TEXT EXTRACTION  (for Knowledge Base indexing)
# ════════════════════════════════════════════════════════════════

SUPPORTED_KB_EXTENSIONS = ["pdf", "docx", "doc", "pptx", "ppt", "xlsx", "xls", "csv", "txt", "md"]

def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """
    Extract plain text from a wide range of file types for KB indexing.

    Supported:
      .pdf   — PyMuPDF (fitz)         — fast, preserves reading order
      .docx  — python-docx            — paragraphs + tables
      .doc   — fallback via docx2txt
      .pptx  — python-pptx            — all slide text frames
      .xlsx / .xls — openpyxl / xlrd  — all sheet cell values
      .csv   — csv stdlib             — all cell values
      .txt / .md — plain UTF-8 decode
    """
    ext = filename.rsplit(".", 1)[-1].lower()

    try:
        # ── PDF ──────────────────────────────────────────────────
        if ext == "pdf":
            import fitz  # PyMuPDF
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            pages = []
            for page in doc:
                pages.append(page.get_text("text"))
            doc.close()
            text = "\n\n".join(pages)

        # ── DOCX ─────────────────────────────────────────────────
        elif ext in ("docx", "doc"):
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text.strip())
            # Also extract table cells
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)
            text = "\n".join(parts)

        # ── PPTX ─────────────────────────────────────────────────
        elif ext in ("pptx", "ppt"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                if slide_parts:
                    parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_parts))
            text = "\n\n".join(parts)

        # ── XLSX / XLS ────────────────────────────────────────────
        elif ext in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                sheet_parts = [f"[Sheet: {sheet.title}]"]
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(
                        str(cell) for cell in row if cell is not None and str(cell).strip()
                    )
                    if row_text:
                        sheet_parts.append(row_text)
                if len(sheet_parts) > 1:
                    parts.append("\n".join(sheet_parts))
            wb.close()
            text = "\n\n".join(parts)

        # ── CSV ───────────────────────────────────────────────────
        elif ext == "csv":
            import csv
            content = file_bytes.decode("utf-8", errors="replace")
            reader = csv.reader(content.splitlines())
            rows = []
            for row in reader:
                row_text = " | ".join(cell.strip() for cell in row if cell.strip())
                if row_text:
                    rows.append(row_text)
            text = "\n".join(rows)

        # ── TXT / MD ──────────────────────────────────────────────
        elif ext in ("txt", "md"):
            text = file_bytes.decode("utf-8", errors="replace")

        else:
            raise ValueError(f"Unsupported file type: .{ext}")

        text = text.strip()
        if not text:
            raise ValueError(f"No text could be extracted from {filename}. "
                             "The file may be empty, scanned-only, or password protected.")
        return text

    except (ValueError, ImportError):
        raise
    except Exception as exc:
        raise RuntimeError(f"Failed to extract text from {filename}: {exc}") from exc


# ════════════════════════════════════════════════════════════════
# SPEECH-TO-TEXT  —  OpenAI Whisper (LOCAL)
# ════════════════════════════════════════════════════════════════

def _get_whisper_model(model_size: str = "base"):
    """Lazy-load Whisper model once. Cached globally."""
    global _whisper_model
    if _whisper_model is None:
        import whisper
        _whisper_model = whisper.load_model(model_size)
    return _whisper_model


def transcribe_audio_bytes(
    audio_bytes: bytes,
    fmt: str = "wav",
    input_lang: str = "auto",
) -> dict:
    """
    Transcribe using OpenAI Whisper locally.

    input_lang: ISO code hint ("auto" = Whisper auto-detects).
                Passing the correct language dramatically improves accuracy
                for non-English speech.

    Returns: {transcribed_text, detected_language}
    """
    try:
        # Convert non-WAV formats
        if fmt not in ("wav", "wave"):
            try:
                from pydub import AudioSegment
                seg = AudioSegment.from_file(io.BytesIO(audio_bytes), format=fmt)
                buf = io.BytesIO()
                seg.export(buf, format="wav")
                audio_bytes = buf.getvalue()
            except Exception as exc:
                raise RuntimeError(f"Audio conversion failed: {exc}")

        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
            tmp.write(audio_bytes)
            tmp_path = tmp.name

        try:
            model = _get_whisper_model("base")

            # Build Whisper transcribe kwargs
            kwargs: dict = {"fp16": False}

            # Pass language hint if user specified one (not auto)
            if input_lang and input_lang not in ("auto", "unknown"):
                whisper_lang = WHISPER_LANG_MAP.get(input_lang, input_lang)
                if whisper_lang:
                    kwargs["language"] = whisper_lang

            result = model.transcribe(tmp_path, **kwargs)
            text = result["text"].strip()
            whisper_detected = result.get("language", "unknown")

            # Normalize detected lang code
            lang_code = _normalize_whisper_lang(whisper_detected)

            return {"transcribed_text": text, "detected_language": lang_code}
        finally:
            try:
                os.unlink(tmp_path)
            except Exception:
                pass

    except Exception as exc:
        raise RuntimeError(f"Transcription failed: {exc}") from exc


def _normalize_whisper_lang(whisper_lang: str) -> str:
    """Normalize Whisper's output language code to our system codes."""
    mapping = {
        "zh": "zh-cn",
        "chinese": "zh-cn",
        "he": "iw",
        "hebrew": "iw",
    }
    return mapping.get(whisper_lang.lower(), whisper_lang)


# ════════════════════════════════════════════════════════════════
# EMBEDDINGS  —  sentence-transformers + FAISS
# ════════════════════════════════════════════════════════════════

_embed_model = None
_faiss_index = None
_faiss_texts: list[str] = []
_faiss_meta:  list[dict] = []   # metadata per chunk: {source, lang}


def _get_embed_model():
    global _embed_model
    if _embed_model is None:
        from sentence_transformers import SentenceTransformer
        _embed_model = SentenceTransformer("all-MiniLM-L6-v2")
    return _embed_model


def _embed(texts: list[str]):
    return _get_embed_model().encode(
        texts, normalize_embeddings=True, show_progress_bar=False
    )


def chunk_text(text: str, size: int = 400, overlap: int = 60) -> list[str]:
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunk = " ".join(words[i: i + size])
        if chunk.strip():
            chunks.append(chunk.strip())
        i += size - overlap
    return chunks


def add_text_to_store(text: str, source: str = "manual", lang: str = "unknown") -> int:
    global _faiss_index, _faiss_texts, _faiss_meta
    try:
        import faiss

        chunks = chunk_text(text)
        if not chunks:
            return 0

        vecs = _embed(chunks).astype("float32")
        dim = vecs.shape[1]

        if _faiss_index is None:
            _faiss_index = faiss.IndexFlatIP(dim)

        _faiss_index.add(vecs)
        _faiss_texts.extend(chunks)
        _faiss_meta.extend([{"source": source, "lang": lang}] * len(chunks))
        return len(chunks)

    except Exception as exc:
        raise RuntimeError(f"Vector store failed: {exc}") from exc


def retrieve_context(query: str, top_k: int = 5) -> str:
    global _faiss_index, _faiss_texts
    if _faiss_index is None or _faiss_index.ntotal == 0:
        return ""
    try:
        import faiss

        q = _embed([query]).astype("float32")
        D, I = _faiss_index.search(q, min(top_k, _faiss_index.ntotal))

        # Only include chunks with similarity score > 0.3
        results = [
            _faiss_texts[i] for score, i in zip(D[0], I[0])
            if i != -1 and score > 0.3
        ]
        return "\n\n---\n\n".join(results)
    except Exception:
        return ""


def get_store_size() -> int:
    return _faiss_index.ntotal if _faiss_index else 0


def clear_vector_store():
    global _faiss_index, _faiss_texts, _faiss_meta
    _faiss_index = None
    _faiss_texts = []
    _faiss_meta  = []


# ════════════════════════════════════════════════════════════════
# GEMINI 2.5 FLASH  —  google-genai SDK
# ════════════════════════════════════════════════════════════════

def ask_gemini(
    question: str,
    api_key: str,
    context: str = "",
    history: Optional[list[dict]] = None,
    response_language: str = "en",
) -> str:
    """
    Gemini 2.5 Flash via new google-genai SDK.
    Model: gemini-2.5-flash
    """
    try:
        from google import genai
        from google.genai import types

        if not question or not question.strip():
            raise ValueError("Question is empty — please record your voice and try again.")

        client = genai.Client(api_key=api_key)

        system_prompt = textwrap.dedent(f"""
            You are SignLingua Assistant — a smart multilingual AI embedded in an OCR
            and voice translation application.

            Rules:
            1. Answer the user's question clearly, accurately and helpfully.
            2. If retrieved context from sign boards or documents is provided, use it
               to ground your answer precisely.
            3. Always respond in the language corresponding to ISO code: {response_language}
            4. Keep answers concise — they will be read aloud via text-to-speech.
            5. If no context is available, answer from your general knowledge.
            6. Never make up facts. If unsure, say so clearly.
        """).strip()

        # Build full conversation history
        contents = []
        for turn in (history or []):
            role = "user" if turn["role"] == "user" else "model"
            if turn["content"].strip():   # skip empty turns
                contents.append(types.Content(
                    role=role,
                    parts=[types.Part(text=turn["content"])]
                ))

        # Add current question (with RAG context if available)
        current_prompt = (
            f"[Relevant context from knowledge base]\n{context}\n\n"
            f"[User question]\n{question}"
            if context else question
        )
        contents.append(types.Content(
            role="user",
            parts=[types.Part(text=current_prompt)]
        ))

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=contents,
            config=types.GenerateContentConfig(
                system_instruction=system_prompt,
                temperature=0.4,
                top_p=0.9,
                max_output_tokens=800,
            ),
        )
        return response.text.strip()

    except ValueError:
        raise
    except Exception as exc:
        raise RuntimeError(f"Gemini API error: {exc}") from exc


# ════════════════════════════════════════════════════════════════
# FULL PIPELINES
# ════════════════════════════════════════════════════════════════

def image_pipeline(
    image_bytes: bytes,
    target_lang: str,
    source_lang: str = "auto",
    add_to_store: bool = True,
) -> dict:
    """
    Image → OCR (with lang hint) → Translate → TTS

    source_lang: hint for Tesseract to improve recognition of non-Latin scripts.
    target_lang: translation output language.
    """
    extracted = extract_text_from_image(image_bytes, hint_lang=source_lang)

    # Translate — use provided source_lang or auto-detect
    src_for_translate = source_lang if source_lang != "auto" else "auto"
    tr = translate_text(extracted, target_lang, source=src_for_translate)

    audio = text_to_speech(tr["translated_text"], target_lang)

    chunks_added = 0
    if add_to_store and extracted:
        chunks_added = add_text_to_store(
            extracted, source="image_ocr", lang=tr["source_language"]
        )

    src = tr["source_language"]
    return {
        "extracted_text":       extracted,
        "source_language":      src,
        "source_language_name": LANG_CODE_TO_NAME.get(src, src.upper()),
        "translated_text":      tr["translated_text"],
        "target_language":      target_lang,
        "audio_bytes":          audio,
        "chunks_indexed":       chunks_added,
    }


def voice_translate_pipeline(
    audio_bytes: bytes,
    target_lang: str,
    input_lang: str = "auto",
) -> dict:
    """
    Live voice → Whisper STT (with lang hint) → Translate → TTS

    input_lang: user-selected input language — passed to Whisper for accuracy.
                "auto" = Whisper detects automatically.
    """
    stt = transcribe_audio_bytes(audio_bytes, fmt="wav", input_lang=input_lang)
    text = stt["transcribed_text"]

    if not text or not text.strip():
        raise ValueError("No speech detected. Speak clearly and try again.")

    det = stt["detected_language"]

    # Use detected or hinted language as source for translation
    src_lang = input_lang if input_lang not in ("auto", "unknown") else det
    tr = translate_text(text, target_lang, source=src_lang)
    audio_out = text_to_speech(tr["translated_text"], target_lang)

    return {
        "transcribed_text":      text,
        "detected_language":     det,
        "detected_language_name": LANG_CODE_TO_NAME.get(det, det.upper()),
        "input_lang_used":       input_lang,
        "translated_text":       tr["translated_text"],
        "target_language":       target_lang,
        "audio_bytes":           audio_out,
    }


def voice_qa_pipeline(
    audio_bytes: bytes,
    api_key: str,
    response_lang: str = "en",
    input_lang: str = "auto",
    history: Optional[list[dict]] = None,
) -> dict:
    """
    Live voice → Whisper STT (with lang hint) → RAG → Gemini 2.5 Flash → TTS

    input_lang: hint for Whisper transcription accuracy.
    response_lang: language for Gemini's reply + TTS output.
    """
    stt = transcribe_audio_bytes(audio_bytes, fmt="wav", input_lang=input_lang)
    question = stt["transcribed_text"]

    if not question or not question.strip():
        raise ValueError("No speech detected. Speak clearly and try again.")

    det = stt["detected_language"]
    context = retrieve_context(question, top_k=5)
    answer = ask_gemini(question, api_key, context, history, response_lang)
    audio_out = text_to_speech(answer, response_lang)

    return {
        "transcribed_question":  question,
        "detected_language":     det,
        "detected_language_name": LANG_CODE_TO_NAME.get(det, det.upper()),
        "context_used":          context,
        "answer_text":           answer,
        "answer_audio_bytes":    audio_out,
    }
